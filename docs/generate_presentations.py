"""
Génère les deux présentations CreditFlow :

  1. CreditFlow-Pitch.pptx                  — argumentaire commercial
  2. CreditFlow-Demo-Fonctionnalites.pptx   — présentation guidée, pas à pas

Usage :  python docs/generate_presentations.py
"""

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Charte graphique (reprend celle de l'application)
# --------------------------------------------------------------------------
BLUE = RGBColor(0x21, 0x52, 0x9C)
BLUE_DARK = RGBColor(0x15, 0x3A, 0x72)
BLUE_LIGHT = RGBColor(0x4B, 0x7A, 0xC8)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
ORANGE = RGBColor(0xE3, 0x74, 0x00)
RED = RGBColor(0xC5, 0x22, 0x1F)
INK = RGBColor(0x1C, 0x23, 0x31)
MUTED = RGBColor(0x5F, 0x6B, 0x7C)
BG = RGBColor(0xF4, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD8, 0xDF, 0xEA)

FONT = "Segoe UI"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_W = SLIDE_W - 2 * MARGIN


# --------------------------------------------------------------------------
# Primitives de mise en page
# --------------------------------------------------------------------------
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=Pt(1)):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = line_w
    box.shadow.inherit = False
    return box


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(0), line_spacing=1.0):
    """runs = [(texte, taille_pt, gras, couleur), ...] — un paragraphe par entrée."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, (content, size, bold, color) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = space_after
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = content
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def wrapped_lines(content: str, width_in: float, size_pt: float) -> int:
    """Nombre de lignes qu'occupera `content` — largeur moyenne d'un glyphe ≈ 0,5 × taille."""
    chars_per_line = max(1, int(width_in * 72 / (size_pt * 0.5)))
    return max(1, math.ceil(len(content) / chars_per_line))


def fit_size(content: str, width_in: float, max_size: float,
             max_lines: int = 1, min_size: float = 13) -> float:
    """Plus grande taille de police tenant en `max_lines` lignes."""
    size = max_size
    while size > min_size and wrapped_lines(content, width_in, size) > max_lines:
        size -= 1
    return size


def page_header(slide, kicker, title, accent=BLUE):
    """Bandeau de titre commun aux slides de contenu."""
    rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=WHITE)
    rect(slide, Emu(0), Emu(0), Inches(0.16), SLIDE_H, fill=accent)
    if kicker:
        text(slide, MARGIN, Inches(0.55), CONTENT_W, Inches(0.3),
             [(kicker.upper(), 12, True, accent)])
    text(slide, MARGIN, Inches(0.9), CONTENT_W, Inches(0.75),
         [(title, 30, True, INK)])
    rect(slide, MARGIN, Inches(1.72), Inches(1.1), Inches(0.045), fill=accent)


def footer(slide, left_text, page_no):
    text(slide, MARGIN, Inches(6.92), Inches(9), Inches(0.3),
         [(left_text, 10, False, MUTED)])
    text(slide, SLIDE_W - MARGIN - Inches(1.2), Inches(6.92), Inches(1.2), Inches(0.3),
         [(str(page_no), 10, True, MUTED)], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, items, size=16, gap=Inches(0.62), bullet_color=BLUE,
            title_size=None):
    """items = [(titre, description|None), ...] — l'avance suit la hauteur réelle du texte."""
    cy = y
    text_w_in = (w - Inches(0.34)) / Inches(1)
    for head, desc in items:
        rect(slide, x, cy + Inches(0.09), Inches(0.13), Inches(0.13),
             fill=bullet_color, shape=MSO_SHAPE.OVAL)
        head_size = title_size or size
        runs = [(head, head_size, True, INK)]
        used = wrapped_lines(head, text_w_in, head_size) * head_size * 1.15
        if desc:
            runs.append((desc, size - 2, False, MUTED))
            used += wrapped_lines(desc, text_w_in, size - 2) * (size - 2) * 1.15 + 4
        text(slide, x + Inches(0.34), cy, w - Inches(0.34),
             Inches(0.5), runs, space_after=Pt(3), line_spacing=1.15)
        cy += max(gap, Inches(used / 72 + 0.16))
    return cy


def stat_cards(slide, y, cards, height=Inches(1.5)):
    """cards = [(valeur, libellé, couleur), ...] — répartis sur toute la largeur.

    La valeur peut être un chiffre court ou un intitulé plus long : la taille de
    police s'adapte pour tenir sur une seule ligne et ne jamais recouvrir le libellé.
    """
    n = len(cards)
    gap = Inches(0.25)
    w = int((CONTENT_W - gap * (n - 1)) / n)
    inner_in = (Emu(w) - Inches(0.5)) / Inches(1)
    for i, (value, label, color) in enumerate(cards):
        x = MARGIN + i * (w + gap)
        rect(slide, x, y, Emu(w), height, fill=BG, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(slide, x, y, Emu(w), Inches(0.075), fill=color)
        value_size = fit_size(value, inner_in, 26, max_lines=1, min_size=14)
        text(slide, x + Inches(0.25), y + Inches(0.36), Emu(w) - Inches(0.5), Inches(0.5),
             [(value, value_size, True, color)])
        text(slide, x + Inches(0.25), y + Inches(0.95), Emu(w) - Inches(0.5),
             height - Inches(1.05),
             [(label, 12, False, MUTED)], line_spacing=1.15)


def panel(slide, x, y, w, h, title, lines, accent=BLUE, mono=False):
    """Encadré titré avec une liste de lignes (texte, gras)."""
    rect(slide, x, y, w, h, fill=BG, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x, y, w, Inches(0.075), fill=accent)
    text(slide, x + Inches(0.3), y + Inches(0.28), w - Inches(0.6), Inches(0.35),
         [(title, 14, True, accent)])
    cy = y + Inches(0.78)
    inner_in = (w - Inches(0.6)) / Inches(1)
    for line, bold in lines:
        n_lines = wrapped_lines(line, inner_in, 13) if line else 1
        h = Inches(n_lines * 13 * 1.2 / 72 + 0.08)
        text(slide, x + Inches(0.3), cy, w - Inches(0.6), h,
             [(line, 13, bold, INK if bold else MUTED)], line_spacing=1.15)
        cy += max(Inches(0.3), h + Inches(0.06))


def notes(slide, content):
    slide.notes_slide.notes_text_frame.text = content


def cover(prs, eyebrow, title, subtitle, meta):
    slide = blank(prs)
    rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=BLUE_DARK)
    rect(slide, Emu(0), Emu(0), Inches(5.6), SLIDE_H, fill=BLUE)
    # bandes décoratives
    rect(slide, SLIDE_W - Inches(3.2), Emu(0), Inches(0.5), SLIDE_H, fill=BLUE)
    rect(slide, SLIDE_W - Inches(2.4), Emu(0), Inches(0.18), SLIDE_H, fill=BLUE_LIGHT)

    text(slide, MARGIN, Inches(2.05), Inches(10), Inches(0.4),
         [(eyebrow.upper(), 13, True, BLUE_LIGHT)])
    text(slide, MARGIN, Inches(2.5), Inches(10.5), Inches(1.5),
         [(title, 54, True, WHITE)], line_spacing=1.0)
    text(slide, MARGIN, Inches(4.05), Inches(9.5), Inches(0.9),
         [(subtitle, 20, False, RGBColor(0xC8, 0xD6, 0xEC))], line_spacing=1.25)
    rect(slide, MARGIN, Inches(5.15), Inches(1.4), Inches(0.05), fill=BLUE_LIGHT)
    text(slide, MARGIN, Inches(5.5), Inches(9), Inches(0.4),
         [(meta, 13, False, RGBColor(0x9F, 0xB4, 0xD6))])
    return slide


def section(prs, number, title, subtitle):
    slide = blank(prs)
    rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=BLUE_DARK)
    text(slide, MARGIN, Inches(2.4), Inches(3), Inches(1.4),
         [(number, 90, True, RGBColor(0x37, 0x63, 0xA8))])
    text(slide, MARGIN + Inches(2.2), Inches(2.75), Inches(9), Inches(0.9),
         [(title, 40, True, WHITE)])
    text(slide, MARGIN + Inches(2.2), Inches(3.85), Inches(9), Inches(0.6),
         [(subtitle, 17, False, RGBColor(0xA8, 0xC0, 0xE2))])
    return slide


def table(slide, x, y, w, headers, rows, col_ratios, row_h=Inches(0.52),
          head_color=BLUE):
    """Tableau simple dessiné à la main (contrôle total du rendu)."""
    widths = [int(w * r) for r in col_ratios]
    # en-tête
    cx = x
    for i, head in enumerate(headers):
        rect(slide, cx, y, Emu(widths[i]), row_h, fill=head_color)
        text(slide, cx + Inches(0.18), y + Inches(0.13), Emu(widths[i]) - Inches(0.3),
             Inches(0.3), [(head, 13, True, WHITE)])
        cx += widths[i]
    # lignes
    cy = y + row_h
    for r, row in enumerate(rows):
        fill = WHITE if r % 2 == 0 else BG
        cx = x
        for i, cell in enumerate(row):
            rect(slide, cx, cy, Emu(widths[i]), row_h, fill=fill, line=BORDER,
                 line_w=Pt(0.75))
            color = INK if i == 0 else MUTED
            bold = i == 0
            if cell.startswith("✓"):
                color, bold = GREEN, True
            elif cell.startswith("✕"):
                color, bold = RED, True
            text(slide, cx + Inches(0.18), cy + Inches(0.14), Emu(widths[i]) - Inches(0.3),
                 Inches(0.32), [(cell, 12.5, bold, color)])
            cx += widths[i]
        cy += row_h
    return cy


# ==========================================================================
# 1. PITCH
# ==========================================================================
def build_pitch(path: Path):
    prs = new_deck()
    p = 0

    cover(prs, "Solution de gestion pour boutique",
          "CreditFlow",
          "Digitaliser la vente à crédit de téléphones et d'ordinateurs :\n"
          "clients, échéances, paiements, relances — même sans réseau.",
          "Présentation commerciale — 14 chantiers livrés, prêt pour un déploiement client")

    # --- Le problème ---------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Le constat", "Aujourd'hui, tout repose sur la mémoire et le cahier", RED)
    bullets(s, MARGIN, Inches(2.15), Inches(6.6), [
        ("Les créances sont suivies à la main",
         "Cahiers, tableurs, notes sur téléphone : l'information est éparpillée."),
        ("Les relances sont envoyées en masse",
         "Un même message générique à tous les clients, chaque début de mois."),
        ("Personne ne sait qui doit quoi, exactement",
         "Le montant restant d'un client demande de reconstituer tout son historique."),
        ("Les retards se découvrent trop tard",
         "Aucune alerte : un impayé se remarque quand le client ne vient plus."),
    ])
    panel(s, MARGIN + Inches(7.1), Inches(2.15), Inches(4.7), Inches(3.3),
          "Le message envoyé chaque mois", [
              ("« Bonjour chers clients, je viens vous", False),
              ("rappeler que le versement mensuel de", False),
              ("votre créance est en cours du 01 au 10", False),
              ("Août 2026... »", False),
              ("", False),
              ("Envoyé à tout le monde, sans montant,", True),
              ("sans échéance, sans suivi.", True),
          ], accent=RED)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Le coût -------------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "L'impact", "Ce que la gestion manuelle coûte réellement", RED)
    stat_cards(s, Inches(2.15), [
        ("Trésorerie", "De l'argent dû qui rentre en retard, ou jamais", RED),
        ("Temps perdu", "Des heures chaque mois à recopier et à recompter", ORANGE),
        ("Erreurs", "Un versement oublié, un montant faux, un litige client", ORANGE),
        ("Aucune visibilité", "Impossible de piloter : combien reste-t-il à récupérer ?", RED),
    ], height=Inches(1.75))
    panel(s, MARGIN, Inches(4.35), CONTENT_W, Inches(1.75),
          "Le vrai problème",
          [("Le commerçant ne manque pas de clients : il manque d'information fiable, "
            "au bon moment.", True),
           ("Chaque créance mal suivie est une vente déjà réalisée dont l'argent "
            "n'est pas encore rentré.", False)],
          accent=RED)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- La solution ---------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "La solution", "CreditFlow : la vente à crédit, entièrement digitalisée")
    text(s, MARGIN, Inches(2.05), CONTENT_W, Inches(0.5),
         [("Une application web unique qui remplace le cahier, le tableur "
           "et les messages envoyés au hasard.", 17, False, MUTED)], line_spacing=1.2)
    cards = [
        ("Enregistrer", "Clients, produits et contrats de crédit en quelques clics", BLUE),
        ("Calculer", "Échéancier, mensualité et date de fin générés automatiquement", BLUE),
        ("Encaisser", "Un versement met tout à jour : échéances, contrat, tableau de bord", GREEN),
        ("Relancer", "Le message personnalisé est prêt à copier pour chaque client", ORANGE),
    ]
    n = len(cards)
    gap = Inches(0.25)
    w = int((CONTENT_W - gap * (n - 1)) / n)
    for i, (title, desc, color) in enumerate(cards):
        x = MARGIN + i * (w + gap)
        rect(s, x, Inches(2.9), Emu(w), Inches(2.5), fill=WHITE, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(2.9), Emu(w), Inches(0.09), fill=color)
        text(s, x + Inches(0.28), Inches(3.25), Emu(w) - Inches(0.56), Inches(0.4),
             [(f"0{i + 1}", 15, True, color)])
        text(s, x + Inches(0.28), Inches(3.7), Emu(w) - Inches(0.56), Inches(0.4),
             [(title, 20, True, INK)])
        text(s, x + Inches(0.28), Inches(4.2), Emu(w) - Inches(0.56), Inches(1.0),
             [(desc, 13, False, MUTED)], line_spacing=1.2)
    panel(s, MARGIN, Inches(5.65), CONTENT_W, Inches(1.0),
          "Principe directeur",
          [("Le commerçant saisit le strict minimum. Le système calcule tout le reste.", True)],
          accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Avant / après -------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Ce qui change", "Avant / après CreditFlow")
    table(s, MARGIN, Inches(2.15), int(CONTENT_W),
          ["", "Avant", "Avec CreditFlow"],
          [
              ["Suivi des créances", "Cahier et mémoire", "✓ Fiche client complète et à jour"],
              ["Calcul des échéances", "À la main, source d'erreurs", "✓ Généré automatiquement"],
              ["Enregistrer un paiement", "Une ligne dans un cahier", "✓ 3 clics, tout se met à jour"],
              ["Connaître le reste dû", "Reconstituer l'historique", "✓ Affiché en permanence"],
              ["Clients en retard", "Découverts trop tard", "✓ Tableau dédié, jours de retard"],
              ["Relances", "Message générique groupé", "✓ Personnalisé, envoi WhatsApp automatique"],
              ["Rapports", "Inexistants", "✓ Export PDF et Excel, taux de défaut inclus"],
              ["Qui a fait quoi", "Aucune trace", "✓ Journal d'audit complet"],
              ["Plusieurs boutiques", "Un cahier par boutique", "✓ Vue consolidée ou par boutique"],
              ["Coupure réseau", "Vente à l'arrêt", "✓ Encaissement hors-ligne, sync automatique"],
              ["Documents remis au client", "Rien d'officiel, ou fait à la main", "✓ Facture et bon de livraison PDF"],
          ],
          [0.26, 0.34, 0.40], row_h=Inches(0.385))
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Le coeur du produit -------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Au cœur du produit", "L'échéancier se calcule tout seul, au franc près")
    text(s, MARGIN, Inches(2.05), Inches(5.6), Inches(1.6),
         [("Le commerçant saisit quatre informations :", 15, False, MUTED),
          ("le prix, l'acompte, le nombre de mensualités et la date de début.",
           15, False, MUTED)], line_spacing=1.25)
    panel(s, MARGIN, Inches(2.75), Inches(5.5), Inches(2.5),
          "Saisie du commerçant", [
              ("Prix de vente", False), ("110 000 FCFA", True),
              ("Acompte", False), ("10 000 FCFA", True),
              ("Mensualités", False), ("3 mois", True),
          ])
    panel(s, MARGIN + Inches(6.0), Inches(2.75), Inches(5.8), Inches(2.5),
          "Calculé par CreditFlow", [
              ("Échéance 1  →  33 333 FCFA", True),
              ("Échéance 2  →  33 333 FCFA", True),
              ("Échéance 3  →  33 334 FCFA", True),
              ("", False),
              ("Total = 100 000 FCFA, exactement.", True),
          ], accent=GREEN)
    panel(s, MARGIN, Inches(5.5), CONTENT_W, Inches(1.05),
          "Pourquoi c'est important",
          [("La dernière échéance absorbe l'arrondi : la somme des échéances est toujours "
            "exactement égale au montant dû. Aucun centime ne se perd, aucun litige possible.",
            False)], accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Relances ------------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Relances", "Chaque client reçoit son message, avec son montant", ORANGE)
    bullets(s, MARGIN, Inches(2.15), Inches(6.3), [
        ("Le tableau des retards est automatique",
         "Jours de retard, montant dû, téléphone : tout est affiché, trié par urgence."),
        ("Un bouton « Générer la relance »",
         "Le message est écrit avec le nom du client et son montant réel."),
        ("Envoi automatique par lot, ou copie manuelle",
         "WhatsApp Cloud API en un clic pour tous les retards, ou copier-coller au cas par cas."),
        ("Chaque envoi est historisé",
         "Date, canal et statut de la relance sont conservés — aucun oubli."),
    ])
    panel(s, MARGIN + Inches(6.8), Inches(2.15), Inches(5.0), Inches(2.9),
          "Message généré", [
              ("Bonjour Ndeye Faye,", True),
              ("Nous vous rappelons que votre", False),
              ("versement mensuel de 39 844 FCFA", False),
              ("est attendu entre le 1 et le 10 du mois.", False),
              ("Merci pour votre confiance.", False),
          ], accent=ORANGE)
    panel(s, MARGIN, Inches(5.35), CONTENT_W, Inches(1.05),
          "Déjà en production",
          [("L'envoi automatique par WhatsApp Cloud API est opérationnel : le commerçant "
            "choisit son canal, la copie manuelle reste disponible en secours.", False)],
          accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Chantiers livrés --------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Journal des livraisons", "Quatorze chantiers, tous fermés et vérifiés", GREEN)
    table(s, MARGIN, Inches(1.95), int(CONTENT_W),
          ["Priorité", "Chantier", "Ce qu'il apporte"],
          [
              ["P0", "Comptes vendeur / caissier", "Rôles distincts : chacun n'accède qu'à ce qui le concerne"],
              ["P0", "Journal d'audit", "Qui a encaissé, annulé ou modifié un prix — et quand"],
              ["P0", "Taux d'intérêt / frais de dossier", "Le prix à crédit reflète le modèle économique réel"],
              ["P1", "Pénalité de retard configurable", "Un levier de recouvrement, pas qu'un indicateur"],
              ["P1", "Relances automatiques", "SMS ou WhatsApp en masse, copie manuelle en secours"],
              ["P1", "Signature électronique", "Pièce d'identité et signature jointes au contrat"],
              ["P1", "Garant / caution", "Sécurise le recouvrement sur les montants élevés"],
              ["P1", "Contrat depuis la fiche client", "Un client, un clic : plus besoin de le rechercher deux fois"],
              ["P2", "Achats fournisseurs", "Le stock reflète les entrées réelles, pas que les ventes"],
              ["P2", "Statistiques avancées", "Taux de défaut par profil, performance par vendeur"],
              ["P2", "Multi-boutiques", "Un tableau de bord consolidé, ou filtré par boutique"],
              ["P2", "Mode hors-ligne", "Un paiement saisi sans réseau se synchronise seul"],
              ["P2", "Facture PDF", "Document officiel prêt à remettre au client à chaque vente"],
              ["P2", "Bon de livraison", "Preuve signée de ce qui a été remis, distincte de la facture"],
          ],
          [0.12, 0.30, 0.58], row_h=Inches(0.32))
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Pilotage ------------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Pilotage", "Le commerçant voit enfin son activité en un coup d'œil")
    stat_cards(s, Inches(2.15), [
        ("Reste à récupérer", "Le total des créances en cours, en temps réel", BLUE),
        ("Encaissé ce mois", "Ce qui est réellement rentré", GREEN),
        ("Clients en retard", "Qui relancer, et de combien de jours", RED),
        ("Prochaines échéances", "Ce qui est attendu dans les 30 jours", ORANGE),
    ], height=Inches(1.7))
    text(s, MARGIN, Inches(4.25), CONTENT_W, Inches(0.4),
         [("Et six rapports exportables en PDF et Excel :", 16, True, INK)])
    bullets(s, MARGIN, Inches(4.8), CONTENT_W, [
        ("Paiements du jour  •  Paiements du mois  •  Clients en retard  •  Créances restantes",
         None),
        ("Taux de défaut par profil client  •  Performance de chaque vendeur (réservé au gérant)",
         None),
    ], size=15, gap=Inches(0.34))
    panel(s, MARGIN, Inches(5.6), CONTENT_W, Inches(0.95),
          "Bénéfice",
          [("Des chiffres fiables pour décider : relancer, accorder un nouveau crédit, "
            "ou réapprovisionner.", False)], accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- UX ------------------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Conçu pour le terrain", "Utilisable par un commerçant, pas par un informaticien")
    bullets(s, MARGIN, Inches(2.2), Inches(5.6), [
        ("Trois clics maximum pour encaisser", "L'opération la plus fréquente est la plus rapide."),
        ("Navigation latérale, toujours visible", "Huit rubriques, pas de menu caché."),
        ("Interface en français, mode clair", "Montants en FCFA, dates au format local."),
        ("Fonctionne sur téléphone et tablette", "Utilisable au comptoir comme en déplacement."),
    ])
    panel(s, MARGIN + Inches(6.3), Inches(2.2), Inches(5.5), Inches(3.4),
          "Encaisser un paiement", [
              ("1.  Ouvrir la liste des retards", True),
              ("2.  Cliquer sur « Encaisser »", True),
              ("3.  Valider le montant proposé", True),
              ("", False),
              ("Le montant de la mensualité est", False),
              ("déjà pré-rempli.", False),
          ], accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Solidité technique --------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Fiabilité", "Une base technique sérieuse, pas un prototype")
    bullets(s, MARGIN, Inches(2.15), Inches(6.2), [
        ("Technologies éprouvées", "Java 21 / Spring Boot, React, PostgreSQL, Docker."),
        ("Accès protégé", "Rôles par utilisateur, session sécurisée par jeton, journal d'audit."),
        ("Données maîtrisées", "Base versionnée par migrations, hébergée chez le client, aucune "
         "dépendance cloud tierce."),
        ("Testé et vérifié", "322 tests backend et 16 tests dédiés à la synchronisation "
         "hors-ligne, tous au vert."),
        ("Sauvegardé automatiquement", "Sauvegarde au démarrage puis toutes les 24 h, "
         "14 jours de rétention, restauration testée."),
    ])
    panel(s, MARGIN + Inches(6.7), Inches(2.15), Inches(5.1), Inches(3.2),
          "Installation", [
              ("Une seule commande :", False),
              ("docker compose up", True),
              ("", False),
              ("Base de données, serveur et", False),
              ("interface démarrent ensemble,", False),
              ("avec un jeu de démonstration.", False),
          ], accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Multi-boutiques -------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Grandir sans changer d'outil", "Une boutique aujourd'hui, plusieurs demain")
    text(s, MARGIN, Inches(2.05), CONTENT_W, Inches(0.5),
         [("Le même outil accompagne l'ouverture d'un second point de vente, sans migration "
           "ni nouvel abonnement.", 15, False, MUTED)], line_spacing=1.2)
    bullets(s, MARGIN, Inches(2.7), Inches(6.2), [
        ("Un utilisateur voit sa boutique, et seulement la sienne",
         "Clients, produits, contrats et paiements restent cloisonnés par boutique."),
        ("Le gérant voit tout, ou une boutique à la fois",
         "Tableau de bord consolidé par défaut, filtre d'un clic sur une boutique précise."),
        ("Les rapports suivent le même principe",
         "Aucune régression sur le fonctionnement mono-boutique existant."),
    ])
    panel(s, MARGIN + Inches(6.6), Inches(2.7), Inches(5.2), Inches(3.0),
          "Vue consolidée du gérant", [
              ("Boutique Centre-ville", True), ("Reste à récupérer : 1 240 000 FCFA", False),
              ("Boutique Marché", True), ("Reste à récupérer : 860 000 FCFA", False),
              ("", False),
              ("Total consolidé  →  2 100 000 FCFA", True),
          ], accent=BLUE)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Mode hors-ligne ---------------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Travailler sans réseau", "Une coupure ne doit jamais coûter un encaissement", GREEN)
    text(s, MARGIN, Inches(2.05), CONTENT_W, Inches(0.5),
         [("Pensé pour le vendeur en tournée de recouvrement, là où la connexion est la "
           "moins fiable.", 15, False, MUTED)], line_spacing=1.2)
    bullets(s, MARGIN, Inches(2.7), Inches(6.2), [
        ("Le paiement se saisit normalement, réseau ou pas",
         "Sans connexion, il apparaît « en attente » au lieu d'être bloqué."),
        ("La synchronisation part seule au retour du réseau",
         "Aucune manipulation : l'application détecte la reconnexion."),
        ("Un conflit est signalé, jamais écrasé en silence",
         "Contrat déjà soldé entre-temps ? Le vendeur est prévenu, rien n'est perdu ni doublé."),
        ("La consultation des contrats déjà ouverts reste disponible",
         "L'application s'installe comme une app mobile, même hors connexion."),
    ])
    panel(s, MARGIN + Inches(6.6), Inches(2.7), Inches(5.2), Inches(3.3),
          "Ce qui protège l'argent encaissé", [
              ("Chaque paiement porte un identifiant unique", False),
              ("généré à la saisie : rejouer la synchronisation", False),
              ("ne peut jamais créer de doublon.", False),
              ("", False),
              ("Vérifié par des tests automatisés dédiés", True),
              ("à ce mécanisme précis.", True),
          ], accent=GREEN)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Avant la mise en service -----------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "En toute franchise", "Ce qui reste à vérifier avant la mise en service", ORANGE)
    bullets(s, MARGIN, Inches(2.15), CONTENT_W, [
        ("Un certificat HTTPS est requis pour le mode hors-ligne complet",
         "Sans lui, l'encaissement hors-ligne fonctionne ; la lecture des fiches sans réseau non."),
        ("Le scénario de conflit mérite un dernier passage terrain",
         "Le mécanisme est testé automatiquement ; une vérification en conditions réelles "
         "reste recommandée avant un déploiement à volume important."),
    ], gap=Inches(0.9))
    panel(s, MARGIN, Inches(4.6), CONTENT_W, Inches(1.3),
          "Pourquoi le dire",
          [("Un produit livré ne veut pas dire un produit qu'on cesse de vérifier. "
            "Ces deux points sont documentés, compris, et n'empêchent pas une mise en "
            "service — ils encadrent la première semaine d'usage.", False)], accent=BLUE)
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Ce qui est livré ----------------------------------------------
    p += 1
    s = blank(prs)
    page_header(s, "Livraison", "Ce qui est disponible dès maintenant", GREEN)
    stat_cards(s, Inches(2.15), [
        ("18", "modules fonctionnels couvrant tout le processus", BLUE),
        ("14 / 14", "chantiers du backlog livrés et vérifiés", GREEN),
        ("338", "tests automatisés au vert (backend + hors-ligne)", GREEN),
        ("1", "commande pour tout démarrer", ORANGE),
    ])
    text(s, MARGIN, Inches(4.15), CONTENT_W, Inches(0.4),
         [("Prêt pour une démonstration immédiate", 17, True, INK)])
    bullets(s, MARGIN, Inches(4.7), CONTENT_W, [
        ("Un jeu de démonstration est créé automatiquement",
         "10 clients, 15 produits, 20 contrats et 30 paiements — dont des retards réalistes."),
        ("Aucune saisie préalable n'est nécessaire pour montrer le produit", None),
    ])
    footer(s, "CreditFlow — Présentation commerciale", p + 1)

    # --- Clôture -------------------------------------------------------
    p += 1
    s = blank(prs)
    rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=BLUE_DARK)
    rect(s, Emu(0), Inches(3.35), SLIDE_W, Inches(0.06), fill=BLUE_LIGHT)
    text(s, MARGIN, Inches(2.2), Inches(11), Inches(1.0),
         [("Reprenez le contrôle de vos créances.", 40, True, WHITE)])
    text(s, MARGIN, Inches(3.7), Inches(10.5), Inches(1.0),
         [("CreditFlow transforme un cahier de créances en un outil de pilotage.",
           19, False, RGBColor(0xC8, 0xD6, 0xEC))], line_spacing=1.25)
    text(s, MARGIN, Inches(5.0), Inches(11), Inches(0.9),
         [("Prochaine étape : une démonstration sur vos propres données.",
           16, True, RGBColor(0x8F, 0xC9, 0xA8))])
    footer_note = "docker compose up  →  http://localhost:3010  →  admin / admin123"
    text(s, MARGIN, Inches(6.3), Inches(11), Inches(0.4),
         [(footer_note, 13, False, RGBColor(0x9F, 0xB4, 0xD6))])

    prs.save(str(path))


# ==========================================================================
# 2. DÉMONSTRATION GUIDÉE
# ==========================================================================
def build_demo(path: Path):
    prs = new_deck()

    def step_slide(num, title, purpose, actions, watch, note, accent=BLUE):
        s = blank(prs)
        page_header(s, f"Étape {num}", title, accent)
        text(s, MARGIN, Inches(2.0), CONTENT_W, Inches(0.45),
             [(purpose, 15, False, MUTED)], line_spacing=1.2)

        # colonne gauche : ce que je fais
        rect(s, MARGIN, Inches(2.65), Inches(6.4), Inches(3.55), fill=WHITE,
             line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, MARGIN, Inches(2.65), Inches(6.4), Inches(0.08), fill=accent)
        text(s, MARGIN + Inches(0.3), Inches(2.95), Inches(5.8), Inches(0.3),
             [("CE QUE JE FAIS À L'ÉCRAN", 11.5, True, accent)])
        cy = Inches(3.45)
        action_gap = min(0.62, 2.6 / max(len(actions), 1))
        for i, action in enumerate(actions, start=1):
            rect(s, MARGIN + Inches(0.3), cy, Inches(0.32), Inches(0.32),
                 fill=accent, shape=MSO_SHAPE.OVAL)
            text(s, MARGIN + Inches(0.3), cy + Inches(0.045), Inches(0.32), Inches(0.3),
                 [(str(i), 11, True, WHITE)], align=PP_ALIGN.CENTER)
            text(s, MARGIN + Inches(0.78), cy + Inches(0.02), Inches(5.35), Inches(0.6),
                 [(action, 13.5, False, INK)], line_spacing=1.15)
            cy += Inches(action_gap)

        # colonne droite : ce qu'il faut montrer
        rect(s, MARGIN + Inches(6.75), Inches(2.65), Inches(5.05), Inches(3.55),
             fill=BG, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, MARGIN + Inches(6.75), Inches(2.65), Inches(5.05), Inches(0.08), fill=GREEN)
        text(s, MARGIN + Inches(7.05), Inches(2.95), Inches(4.45), Inches(0.3),
             [("CE QU'IL FAUT SOULIGNER", 11.5, True, GREEN)])
        cy = Inches(3.45)
        watch_gap = min(0.68, 2.6 / max(len(watch), 1))
        for item in watch:
            rect(s, MARGIN + Inches(7.05), cy + Inches(0.08), Inches(0.12), Inches(0.12),
                 fill=GREEN, shape=MSO_SHAPE.OVAL)
            text(s, MARGIN + Inches(7.35), cy, Inches(4.15), Inches(0.7),
                 [(item, 13, False, INK)], line_spacing=1.15)
            cy += Inches(watch_gap)

        # 3 diapositives précèdent les étapes : couverture, préparation, sommaire
        footer(s, "CreditFlow — Démonstration guidée", num + 3)
        notes(s, note)
        return s

    # --- Couverture ----------------------------------------------------
    cover(prs, "Présentation guidée",
          "CreditFlow en démonstration",
          "Parcours complet en 18 étapes : du client à la relance,\n"
          "multi-boutiques et hors-ligne compris.",
          "Durée indicative : 25 à 30 minutes")

    # --- Avant de commencer --------------------------------------------
    s = blank(prs)
    page_header(s, "Préparation", "Avant de commencer")
    panel(s, MARGIN, Inches(2.15), Inches(5.7), Inches(2.5),
          "Accès", [
              ("Application", False), ("http://localhost:3010", True),
              ("Identifiant", False), ("admin  /  admin123", True),
          ])
    panel(s, MARGIN + Inches(6.1), Inches(2.15), Inches(5.7), Inches(2.5),
          "Jeu de démonstration (créé automatiquement)", [
              ("10 clients   •   15 produits", True),
              ("20 contrats de crédit", True),
              ("30 paiements déjà enregistrés", True),
              ("9 clients en situation de retard", True),
          ], accent=GREEN)
    panel(s, MARGIN, Inches(4.9), CONTENT_W, Inches(1.5),
          "Conseil de démonstration",
          [("Les données de retard sont réelles et calculées à la date du jour : "
            "les chiffres affichés à l'écran évoluent donc naturellement.", False),
           ("Suivre l'ordre des étapes : chaque écran prépare le suivant.", True)],
          accent=ORANGE)
    footer(s, "CreditFlow — Démonstration guidée", 2)

    # --- Le parcours ---------------------------------------------------
    s = blank(prs)
    page_header(s, "Sommaire", "Le parcours en un coup d'œil")
    left = [
        "1.  Connexion", "2.  Tableau de bord", "3.  Boutiques et vue consolidée",
        "4.  Liste des clients", "5.  Fiche client complète",
        "6.  Catalogue produits", "7.  Achats fournisseurs et stock",
        "8.  Créer une vente à crédit", "9.  Le contrat et son échéancier",
    ]
    right = [
        "10. Encaisser un paiement", "11. Les mises à jour automatiques",
        "12. Mode hors-ligne", "13. Suivi des échéances",
        "14. Clients en retard et relance", "15. Rapports et exports",
        "16. Recherche globale", "17. Utilisateurs et rôles", "18. API et documentation",
    ]
    for col, items in ((MARGIN, left), (MARGIN + Inches(6.1), right)):
        cy = Inches(2.05)
        for item in items:
            text(s, col, cy, Inches(5.6), Inches(0.38),
                 [(item, 14, False, INK)])
            rect(s, col, cy + Inches(0.38), Inches(5.4), Inches(0.012), fill=BORDER)
            cy += Inches(0.52)
    footer(s, "CreditFlow — Démonstration guidée", 3)

    # --- Étapes --------------------------------------------------------
    step_slide(
        1, "Connexion",
        "Un compte administrateur créé automatiquement, des comptes vendeurs à la demande.",
        ["Ouvrir http://localhost:3010",
         "Saisir admin / admin123",
         "Cliquer sur « Se connecter »"],
        ["L'accès est protégé : sans connexion, aucune donnée n'est accessible.",
         "Le compte administrateur est créé tout seul, aucune installation manuelle.",
         "La session reste ouverte, le commerçant ne se reconnecte pas sans arrêt."],
        "Insister sur la simplicité du premier démarrage : un seul compte suffit "
        "pour commencer. La création de comptes vendeurs avec des droits limités "
        "est montrée à l'étape 17. Le mot de passe est modifiable dans le fichier .env.")

    step_slide(
        2, "Tableau de bord",
        "La première chose que voit le commerçant le matin : l'état de son activité.",
        ["Parcourir les sept indicateurs du haut",
         "Descendre sur « Clients en retard »",
         "Montrer « Paiements du jour » et « Prochaines échéances »",
         "Cliquer sur une carte pour rebondir vers le détail"],
        ["Nombre de clients, ventes à crédit, montant restant à récupérer.",
         "Montant encaissé ce mois-ci et paiements du jour.",
         "Les clients en retard remontent directement, avec le nombre de jours.",
         "Deux boutons d'action rapide : encaisser et créer une vente."],
        "C'est LA slide qui vend le produit. Tout ce que le commerçant cherchait dans "
        "son cahier est ici, sans rien chercher. Les cartes sont cliquables : "
        "chaque chiffre mène à l'écran correspondant.")

    step_slide(
        3, "Boutiques et vue consolidée",
        "Le tableau de bord change de nature dès qu'une deuxième boutique existe.",
        ["Ouvrir le sélecteur de boutique dans la barre supérieure",
         "Basculer entre « Toutes les boutiques » et une boutique précise",
         "Montrer que les chiffres du tableau de bord se recalculent",
         "Ouvrir « Boutiques » dans le menu (visible ADMIN uniquement)"],
        ["Invisible si une seule boutique existe : aucune complexité imposée au démarrage.",
         "Chaque utilisateur ne voit que les boutiques auxquelles il est rattaché.",
         "La consolidation ne casse aucun rapport ni export existant."],
        "À ne montrer que si le prospect gère ou envisage plusieurs points de vente. "
        "Sinon, passer directement à l'étape suivante : le sélecteur reste discret "
        "pour un commerçant mono-boutique.",
        accent=BLUE_LIGHT)

    step_slide(
        4, "Liste des clients",
        "Retrouver n'importe quel client en une seconde.",
        ["Ouvrir « Clients » dans le menu latéral",
         "Taper un nom dans la recherche (ex. « Diallo »)",
         "Puis chercher par numéro de téléphone",
         "Cliquer sur « Nouveau client » pour montrer le formulaire"],
        ["La recherche fonctionne sur le nom, le téléphone, la CNI et la profession.",
         "Les résultats s'affichent au fur et à mesure de la frappe.",
         "Le formulaire refuse un téléphone ou une CNI déjà utilisés."],
        "Montrer que la recherche accepte un nom partiel. Insister sur le contrôle "
        "des doublons : impossible d'enregistrer deux fois le même client, "
        "ce qui évite les créances éclatées sur deux fiches.")

    step_slide(
        5, "Fiche client complète",
        "Toute la relation commerciale d'un client sur un seul écran.",
        ["Cliquer sur l'icône « voir » d'un client",
         "Montrer le bloc « Situation financière »",
         "Faire défiler l'historique des achats",
         "Puis l'historique des paiements",
         "Cliquer sur « Nouveau contrat » pour vendre à ce client, sans le rechercher à nouveau"],
        ["Total acheté, total payé, reste à payer, échéances en retard.",
         "Tous les contrats du client, avec leur statut.",
         "Tous les versements, avec date, mode et référence.",
         "Une photo peut être ajoutée pour identifier le client.",
         "Le client est déjà présélectionné dans le formulaire de vente qui s'ouvre."],
        "C'est la réponse à la question la plus fréquente du commerçant : "
        "« Combien me doit ce client ? ». Avant, il fallait tout reconstituer. "
        "Montrer le bouton « Nouveau contrat » : c'est le raccourci le plus utilisé "
        "au quotidien, il évite d'aller chercher le client une seconde fois dans "
        "l'écran des ventes. Cliquer sur un contrat existant pour montrer qu'on "
        "rebondit aussi vers son détail.")

    step_slide(
        6, "Catalogue produits",
        "Les téléphones et ordinateurs vendus, avec leurs deux prix.",
        ["Ouvrir « Produits »",
         "Filtrer par catégorie",
         "Montrer un produit avec prix comptant et prix à crédit",
         "Signaler un produit en rupture"],
        ["Deux prix distincts : comptant et à crédit.",
         "Le stock se décrémente automatiquement à chaque vente.",
         "Un produit dont le stock tombe à zéro passe en rupture tout seul."],
        "Le double prix est le cœur du modèle économique de la boutique : "
        "le crédit se paie. Le passage automatique en rupture évite de vendre "
        "un produit qu'on n'a plus.")

    step_slide(
        7, "Achats fournisseurs et stock",
        "Le stock ne se corrige plus à la main : il se reconstitue par une vraie réception.",
        ["Ouvrir « Fournisseurs », montrer la fiche d'un fournisseur",
         "Ouvrir « Achats », créer une réception de stock",
         "Choisir un ou plusieurs produits et leur quantité reçue",
         "Valider, puis revenir sur la fiche produit pour voir le stock augmenté"],
        ["Le stock reflète maintenant les entrées réelles, pas seulement les ventes.",
         "L'historique des mouvements distingue clairement entrée (achat) et sortie (vente).",
         "Aucune correction manuelle du stock n'est plus nécessaire."],
        "Slide pour un commerçant qui gère aussi ses approvisionnements, pas seulement "
        "sa caisse. À raccourcir si l'interlocuteur ne s'occupe que de la vente.")

    step_slide(
        8, "Créer une vente à crédit",
        "L'écran le plus important : c'est ici que naît la créance.",
        ["Cliquer sur « Nouvelle vente »",
         "Choisir un client, puis un produit",
         "Saisir l'acompte et le nombre de mensualités",
         "Renseigner un garant si le montant ou le profil client le justifie",
         "Observer l'échéancier qui se calcule à droite, en direct"],
        ["Le prix se pré-remplit avec le prix à crédit du produit.",
         "L'échéancier apparaît AVANT d'enregistrer : rien n'est une surprise.",
         "Montant à financer, mensualité et date de fin sont calculés seuls.",
         "Le garant est optionnel : nom, téléphone et adresse d'un tiers responsable.",
         "Le commerçant peut ajuster et voir l'effet immédiatement."],
        "MOMENT CLÉ DE LA DÉMO. Faire varier le nombre de mensualités devant le public "
        "et montrer que l'échéancier se recalcule instantanément. Ajouter un garant "
        "sur un montant élevé pour illustrer la sécurisation du recouvrement.",
        accent=GREEN)

    step_slide(
        9, "Le contrat et son échéancier",
        "Un contrat = un client, un produit, un échéancier daté — et ses preuves.",
        ["Ouvrir le contrat qui vient d'être créé",
         "Montrer le résumé et la barre de progression",
         "Télécharger la facture PDF, puis le bon de livraison",
         "Joindre une pièce d'identité scannée ou une signature tactile",
         "Parcourir le tableau des échéances",
         "Montrer une échéance en retard, en rouge"],
        ["Chaque échéance a son numéro, sa date, son montant et son statut.",
         "La somme des échéances est exactement égale au montant dû.",
         "Facture, bon de livraison et reçu de paiement sont trois documents distincts, "
         "chacun avec sa propre référence — jamais de doublon entre eux.",
         "Le bon de livraison porte une zone de signature, prêt à être imprimé.",
         "La pièce jointe est vérifiée sur son contenu réel, pas seulement son extension.",
         "Le retard est calculé en jours, automatiquement.",
         "Une référence unique identifie le contrat (ex. VC-2026-00021)."],
        "Expliquer l'arrondi : sur 100 000 FCFA en 3 fois, on obtient "
        "33 333 + 33 333 + 33 334. La dernière échéance absorbe le reliquat, "
        "donc le compte tombe juste au franc près. C'est ce qui évite les litiges. "
        "Ouvrir réellement la facture et le bon de livraison téléchargés : "
        "montrer qu'ils ne se ressemblent pas, chacun a un objectif précis.")

    step_slide(
        10, "Encaisser un paiement",
        "L'opération la plus fréquente doit être la plus rapide : trois clics.",
        ["Depuis le contrat, cliquer sur « Encaisser »",
         "Constater que le montant est déjà pré-rempli",
         "Choisir le mode de paiement",
         "Valider"],
        ["Le montant proposé est celui de la mensualité : souvent rien à saisir.",
         "Modes disponibles : espèces, Mobile Money, virement, chèque, carte.",
         "Une référence et une observation peuvent être ajoutées.",
         "Le même bouton existe sur le tableau de bord et la liste des retards."],
        "Chronométrer devant le public si possible. C'est l'argument UX principal : "
        "le commerçant encaisse pendant que le client est devant lui, "
        "sans le faire attendre.",
        accent=GREEN)

    step_slide(
        11, "Ce que le paiement met à jour, tout seul",
        "Un seul geste, quatre mises à jour simultanées.",
        ["Revenir sur l'échéancier du contrat",
         "Montrer l'échéance passée à « Payée »",
         "Montrer le reste à payer qui a diminué",
         "Retourner au tableau de bord"],
        ["Le versement solde d'abord les échéances les plus anciennes.",
         "Une échéance partiellement réglée passe en « Partiel ».",
         "Le contrat passe en « Soldé » dès que le reste atteint zéro.",
         "Le tableau de bord reflète immédiatement l'encaissement."],
        "Point technique à vulgariser : la règle d'imputation est FIFO, "
        "on solde toujours le retard le plus ancien en premier. "
        "Exemple concret : un versement de 50 000 sur des échéances de 33 333 "
        "solde la première et règle 16 667 sur la deuxième.",
        accent=GREEN)

    step_slide(
        12, "Mode hors-ligne",
        "La panne réseau ne doit plus jamais interrompre un encaissement.",
        ["Couper le réseau (mode avion, ou DevTools « Offline »)",
         "Encaisser un paiement normalement : il passe « en attente »",
         "Réactiver le réseau",
         "Montrer la synchronisation automatique, sans action du vendeur"],
        ["Le paiement hors-ligne est visible localement avec son statut « en attente ».",
         "La synchronisation part seule dès que la connexion revient.",
         "Un conflit (contrat soldé entre-temps) est signalé, jamais écrasé en silence.",
         "Les fiches déjà consultées restent lisibles sans connexion."],
        "Étape à fort impact pour un vendeur en tournée de recouvrement. Couper "
        "vraiment le réseau devant le public : l'effet de la synchronisation "
        "automatique au retour marque les esprits. Préciser qu'un certificat "
        "HTTPS est nécessaire en production pour bénéficier de tout le mécanisme.",
        accent=GREEN)

    step_slide(
        13, "Suivi des échéances",
        "Toutes les échéances de la boutique, filtrables.",
        ["Ouvrir « Échéances »",
         "Activer le filtre « En retard »",
         "Filtrer sur une période",
         "Rechercher un client précis"],
        ["Colonnes : client, produit, date prévue, montant, statut, retard.",
         "Le statut se calcule en direct, il n'y a rien à rafraîchir.",
         "Chaque ligne permet d'encaisser directement.",
         "Les filtres se combinent librement."],
        "Insister : les retards ne sont pas stockés, ils sont recalculés à chaque "
        "affichage à partir de la date du jour. Il n'y a donc aucun risque "
        "de voir un chiffre périmé.")

    step_slide(
        14, "Clients en retard et relance",
        "La fonction qui remplace le message groupé du début de mois.",
        ["Ouvrir « Relances »",
         "Parcourir le tableau, trié par ancienneté du retard",
         "Cliquer sur « Générer la relance » pour un client",
         "Montrer l'envoi automatique en masse, puis la copie manuelle en secours"],
        ["Jours de retard, montant dû, téléphone : tout est là.",
         "Le message est personnalisé avec le nom et le montant réel du client.",
         "L'envoi par WhatsApp Cloud API est automatique, individuel ou en masse.",
         "Chaque relance envoyée est historisée avec son statut."],
        "Comparer avec l'existant : aujourd'hui un message unique part à tous, "
        "sans montant. Ici chaque client reçoit SON montant, par le canal "
        "automatique — la copie manuelle reste disponible si le canal échoue.",
        accent=ORANGE)

    step_slide(
        15, "Rapports et exports",
        "De quoi rendre des comptes, et archiver.",
        ["Ouvrir « Rapports »",
         "Basculer entre les six rapports disponibles",
         "Montrer le taux de défaut par profession, puis la performance vendeur",
         "Télécharger en PDF, puis en Excel"],
        ["Paiements du jour, paiements du mois, clients en retard, créances restantes.",
         "Taux de défaut par profil client, filtrable par tranche de montant.",
         "Performance par vendeur — visible uniquement par le gérant.",
         "Le PDF est prêt à imprimer ou à archiver, l'Excel à retravailler librement."],
        "Ouvrir réellement un PDF téléchargé pendant la démo : l'effet est fort. "
        "Le rapport « Performance vendeur » est celui qui intéresse le plus "
        "un gérant à plusieurs employés.")

    step_slide(
        16, "Recherche globale",
        "Une seule barre de recherche pour tout retrouver.",
        ["Cliquer dans la barre de recherche en haut",
         "Taper un nom de client",
         "Puis un numéro de téléphone",
         "Puis le nom ou le téléphone d'un garant"],
        ["Cherche simultanément dans les clients, les produits, les contrats et les garants.",
         "Les résultats sont regroupés par catégorie.",
         "Un clic mène directement à la fiche concernée."],
        "Cas d'usage réel : un client appelle et donne seulement son prénom "
        "ou son numéro. En deux secondes le commerçant a toute sa situation.")

    step_slide(
        17, "Utilisateurs et rôles",
        "Le gérant garde la main sur qui peut faire quoi.",
        ["Ouvrir « Utilisateurs » (visible ADMIN uniquement)",
         "Créer un compte vendeur avec un mot de passe temporaire",
         "Montrer les actions masquées pour ce rôle (suppression, prix)",
         "Ouvrir l'historique d'une fiche pour montrer le journal d'audit"],
        ["Un vendeur peut vendre et encaisser, pas supprimer ni changer un prix.",
         "Chaque action sensible est journalisée : qui, quoi, quand.",
         "Un compte désactivé perd l'accès sans effacer son historique."],
        "Slide de confiance pour un gérant qui emploie plusieurs personnes : "
        "il garde le contrôle des opérations sensibles et une trace de tout, "
        "sans avoir à surveiller chaque geste en personne.",
        accent=BLUE_LIGHT)

    step_slide(
        18, "API et documentation",
        "Pour la suite : le produit est ouvert et documenté.",
        ["Ouvrir http://localhost:8080/swagger-ui.html",
         "Parcourir les rubriques",
         "Déplier une opération",
         "Montrer qu'on peut la tester en ligne"],
        ["Toutes les fonctions sont accessibles par API REST documentée.",
         "Une caisse, une comptabilité ou une application mobile peuvent s'y brancher.",
         "C'est la même base qui porte déjà le multi-boutiques et le hors-ligne."],
        "Slide destinée à un interlocuteur technique ou à un investisseur. "
        "Devant un commerçant seul, on peut la passer rapidement. "
        "Le message : le produit n'est pas une impasse technique.",
        accent=BLUE_LIGHT)

    # --- Récapitulatif -------------------------------------------------
    s = blank(prs)
    page_header(s, "Conclusion", "Ce que la démonstration vient de montrer", GREEN)
    stat_cards(s, Inches(2.1), [
        ("Tout est relié", "Une saisie met à jour tous les écrans", BLUE),
        ("Rien à calculer", "Échéances et retards sont automatiques", GREEN),
        ("3 clics", "Pour encaisser un versement", GREEN),
        ("Prêt à l'emploi", "Une commande pour tout démarrer", ORANGE),
    ], height=Inches(1.6))
    panel(s, MARGIN, Inches(4.05), Inches(5.7), Inches(2.2),
          "Le gain pour le commerçant", [
              ("Il sait en permanence qui lui doit quoi.", True),
              ("Il relance les bons clients, avec le bon montant.", True),
              ("Il encaisse sans faire attendre le client.", True),
              ("Il pilote sa trésorerie avec des chiffres fiables.", True),
          ], accent=GREEN)
    panel(s, MARGIN + Inches(6.1), Inches(4.05), Inches(5.7), Inches(2.2),
          "Questions fréquentes", [
              ("« Et si je me trompe de montant ? »", True),
              ("Un versement s'annule, tout est recalculé.", False),
              ("« Et une coupure réseau en pleine vente ? »", True),
              ("Le paiement reste en attente et se", False),
              ("synchronise seul, sans jamais se dupliquer.", False),
          ], accent=ORANGE)
    footer(s, "CreditFlow — Démonstration guidée", 22)

    # --- Clôture -------------------------------------------------------
    s = blank(prs)
    rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=BLUE_DARK)
    rect(s, Emu(0), Inches(3.3), SLIDE_W, Inches(0.06), fill=BLUE_LIGHT)
    text(s, MARGIN, Inches(2.3), Inches(11), Inches(1.0),
         [("Merci — place à vos questions", 40, True, WHITE)])
    text(s, MARGIN, Inches(3.65), Inches(10.5), Inches(0.8),
         [("La démonstration peut être rejouée à tout moment sur vos propres données.",
           18, False, RGBColor(0xC8, 0xD6, 0xEC))], line_spacing=1.25)
    text(s, MARGIN, Inches(4.9), Inches(11), Inches(1.2),
         [("Application       http://localhost:3010", 15, False, RGBColor(0x9F, 0xB4, 0xD6)),
          ("Documentation  http://localhost:8080/swagger-ui.html", 15, False,
           RGBColor(0x9F, 0xB4, 0xD6)),
          ("Démarrage        docker compose up", 15, False, RGBColor(0x9F, 0xB4, 0xD6))],
         space_after=Pt(8))

    prs.save(str(path))


if __name__ == "__main__":
    out = Path(__file__).resolve().parent
    build_pitch(out / "CreditFlow-Pitch.pptx")
    build_demo(out / "CreditFlow-Demo-Fonctionnalites.pptx")
    print("Presentations generees dans", out)
